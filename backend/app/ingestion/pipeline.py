import os
from typing import Dict, List

# Text extractors


def _read_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _read_pdf(path: str) -> str:
    try:
        from pdfminer.high_level import extract_text  # type: ignore
    except Exception:
        return ""
    try:
        return extract_text(path) or ""
    except Exception:
        return ""


def _read_docx(path: str) -> str:
    try:
        from docx import Document  # type: ignore
    except Exception:
        return ""
    try:
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except Exception:
        return ""


def _read_pptx(path: str) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return ""
    try:
        prs = Presentation(path)
        texts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    texts.append(shape.text_frame.text)
        return "\n".join(t for t in texts if t)
    except Exception:
        return ""


def _extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext in {".txt", ".md"}:
        return _read_txt(file_path)
    if ext == ".pdf":
        return _read_pdf(file_path)
    if ext in {".docx", ".doc"}:  # .doc fallback will likely fail; best-effort
        return _read_docx(file_path)
    if ext in {".pptx"}:
        return _read_pptx(file_path)
    return ""


def _chunk_text(text: str, max_chars: int = 1000, overlap: int = 100) -> List[str]:
    text = (text or "").strip()
    if not text:
        return []
    chunks: List[str] = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + max_chars, n)
        chunk = text[start:end]
        chunks.append(chunk)
        if end == n:
            break
        start = end - overlap if end - overlap > start else end
    return [c.strip() for c in chunks if c.strip()]


# Public API
async def process_file_async(file_path: str, course_id: str) -> Dict:
    from ..services.qdrant import upsert_texts  # local import to avoid heavy deps at import time

    ext = os.path.splitext(file_path)[1].lower()
    raw = _extract_text(file_path)
    chunks = _chunk_text(raw)

    result: Dict = {
        "course_id": course_id,
        "file_path": file_path,
        "ext": ext,
        "chunks": len(chunks),
    }

    if not chunks:
        result.update({"status": "no-text"})
        return result

    try:
        upserted = await upsert_texts(
            course_id, chunks, metadata={"source": os.path.basename(file_path), "ext": ext}
        )
        result.update({"status": "ingested", "vector": upserted})
    except Exception as e:
        # Do not crash upload if embeddings/Qdrant are unavailable
        result.update({
            "status": "vectorization_failed",
            "error": str(e),
        })
    return result


def process_file(file_path: str, course_id: str) -> Dict:
    """Deprecated sync shim left for compatibility. Will only summarize basic info."""
    ext = os.path.splitext(file_path)[1].lower()
    return {
        "course_id": course_id,
        "file_path": file_path,
        "ext": ext,
        "chunks": 0,
        "status": "ingested (stub)",
    }
